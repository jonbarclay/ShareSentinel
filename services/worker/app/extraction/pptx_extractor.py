"""PowerPoint (.pptx) text extraction using python-pptx."""

from pathlib import Path

from pptx import Presentation

from .base import BaseExtractor, ExtractionResult


class PptxExtractor(BaseExtractor):
    """Extract text from PowerPoint presentations including slides, tables, and speaker notes."""

    def extract(self, file_path: Path, file_size: int) -> ExtractionResult:
        try:
            prs = Presentation(str(file_path))

            metadata: dict = {
                "slide_count": len(prs.slides),
            }
            # Extract core properties if available
            try:
                props = prs.core_properties
                metadata["title"] = props.title or ""
                metadata["author"] = props.author or ""
                metadata["subject"] = props.subject or ""
                metadata["keywords"] = props.keywords or ""
            except Exception:
                pass  # Properties might not be accessible

            parts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts: list[str] = []

                # Slide text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_parts.append(text)

                    # Table content
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_parts.append(" | ".join(cells))

                # Speaker notes
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                slide_header = f"--- Slide {slide_num} ---"
                if slide_parts:
                    parts.append(slide_header + "\n" + "\n".join(slide_parts))

                if notes_text:
                    parts.append(
                        f"--- Slide {slide_num} Speaker Notes ---\n{notes_text}"
                    )

            full_text = "\n\n".join(parts)

            if len(full_text.strip()) < self.MIN_MEANINGFUL_TEXT:
                return ExtractionResult(
                    success=False,
                    extraction_method="pptx_text",
                    metadata=metadata,
                    error="Presentation appears to be empty or contains only non-text elements",
                )

            text, was_sampled, description = self._truncate_if_needed(
                full_text,
                f"First ~100KB of text from a {metadata['slide_count']}-slide presentation",
            )

            return ExtractionResult(
                success=True,
                extraction_method="pptx_text",
                text_content=text,
                metadata=metadata,
                was_sampled=was_sampled,
                sampling_description=description,
                content_length=len(text),
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                extraction_method="pptx_text",
                error=str(e),
            )
